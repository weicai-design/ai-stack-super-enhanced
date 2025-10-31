"""
Office document handlers for .docx, .xlsx and .pptx files.

Provides: parse_docx, parse_xlsx and parse_pptx.

All functions return a list of chunk dicts with keys:
    text, source, part, start_offset, end_offset, checksum
"""

from __future__ import annotations

import hashlib
import os
from typing import Dict, List


def _file_checksum(path: str) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        while True:
            chunk = f.read(8192)
            if not chunk:
                break
            h.update(chunk)
    return h.hexdigest()


def parse_docx(path: str) -> List[Dict]:
    if not os.path.exists(path):
        raise FileNotFoundError(path)

    try:
        import docx  # python-docx
    except Exception as e:
        raise ImportError("python-docx required") from e

    doc = docx.Document(path)
    checksum = _file_checksum(path)
    chunks = []
    for i, para in enumerate(doc.paragraphs):
        text = para.text or ""
        chunks.append(
            {
                "text": text,
                "source": os.path.abspath(path),
                "part": "paragraph",
                "index": i,
                "start_offset": 0,
                "end_offset": len(text),
                "checksum": checksum,
            }
        )
    return chunks


def parse_xlsx(path: str) -> List[Dict]:
    if not os.path.exists(path):
        raise FileNotFoundError(path)

    try:
        import openpyxl
    except Exception as e:
        raise ImportError("openpyxl is required to parse .xlsx files") from e

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    checksum = _file_checksum(path)
    chunks = []
    for sheet in wb.worksheets:
        for r_idx, row in enumerate(sheet.iter_rows(values_only=True)):
            text = "\t".join([str(c) if c is not None else "" for c in row])
            chunks.append(
                {
                    "text": text,
                    "source": os.path.abspath(path),
                    "part": sheet.title,
                    "index": r_idx,
                    "start_offset": 0,
                    "end_offset": len(text),
                    "checksum": checksum,
                }
            )
    return chunks


def parse_pptx(path: str) -> List[Dict]:
    if not os.path.exists(path):
        raise FileNotFoundError(path)

    try:
        from pptx import Presentation
    except Exception as e:
        raise ImportError("python-pptx required") from e

    prs = Presentation(path)
    checksum = _file_checksum(path)
    chunks = []
    for s_idx, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        text = "\n".join([t for t in texts if t])
        chunks.append(
            {
                "text": text,
                "source": os.path.abspath(path),
                "part": "slide",
                "index": s_idx,
                "start_offset": 0,
                "end_offset": len(text),
                "checksum": checksum,
            }
        )
    return chunks
