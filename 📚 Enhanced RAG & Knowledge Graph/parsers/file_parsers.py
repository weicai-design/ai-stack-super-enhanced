from __future__ import annotations

import io
from pathlib import Path
from typing import Any, Callable, Dict, List


def _split_chunks(text: str, max_len: int = 1200, overlap: int = 100) -> List[str]:
    if not text:
        return []
    chunks = []
    i = 0
    n = len(text)
    while i < n:
        j = min(i + max_len, n)
        chunk = text[i:j]
        chunks.append(chunk)
        if j == n:
            break
        i = max(0, j - overlap)
    return chunks


def _read_text_fallback(p: Path) -> str:
    try:
        return p.read_text(encoding="utf-8")
    except Exception:
        try:
            import chardet  # type: ignore

            raw = p.read_bytes()
            enc = chardet.detect(raw).get("encoding") or "utf-8"
            return raw.decode(enc, errors="ignore")
        except Exception:
            return ""


def parse_txt_md(p: Path) -> str:
    return _read_text_fallback(p)


def parse_pdf(p: Path) -> str:
    try:
        import fitz  # PyMuPDF type: ignore

        doc = fitz.open(p)
        texts = []
        for page in doc:
            texts.append(page.get_text())
        return "\n".join(texts)
    except Exception:
        try:
            import pdfplumber  # type: ignore

            texts = []
            with pdfplumber.open(p) as pdf:
                for page in pdf.pages:
                    texts.append(page.extract_text() or "")
            return "\n".join(texts)
        except Exception:
            return _read_text_fallback(p)


def parse_docx(p: Path) -> str:
    try:
        import docx  # python-docx

        d = docx.Document(str(p))
        return "\n".join([para.text for para in d.paragraphs])
    except Exception:
        return _read_text_fallback(p)


def parse_pptx(p: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore

        prs = Presentation(str(p))
        texts = []
        for slide in prs.slides:
            for shp in slide.shapes:
                if hasattr(shp, "text"):
                    texts.append(shp.text)
        return "\n".join(texts)
    except Exception:
        return _read_text_fallback(p)


def parse_xlsx(p: Path, max_rows: int = 2000) -> str:
    try:
        from openpyxl import load_workbook  # type: ignore

        wb = load_workbook(filename=str(p), read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            cnt = 0
            for row in ws.iter_rows(values_only=True):
                parts.append("\t".join("" if v is None else str(v) for v in row))
                cnt += 1
                if cnt >= max_rows:
                    break
        return "\n".join(parts)
    except Exception:
        return _read_text_fallback(p)


def parse_epub(p: Path) -> str:
    try:
        from bs4 import BeautifulSoup  # type: ignore
        from ebooklib import epub  # type: ignore

        book = epub.read_epub(str(p))
        texts = []
        for item in book.get_items():
            if item.get_type() == 9:  # DOCUMENT
                soup = BeautifulSoup(item.get_body_content(), "html.parser")
                texts.append(soup.get_text(separator=" "))
        return "\n".join(texts)
    except Exception:
        return _read_text_fallback(p)


def parse_csv(p: Path, max_rows: int = 5000) -> str:
    try:
        import pandas as pd  # type: ignore

        df = pd.read_csv(p, nrows=max_rows)
        # 以“制表符分隔”文本化
        buf = io.StringIO()
        df.to_csv(buf, sep="\t", index=False)
        return buf.getvalue()
    except Exception:
        return _read_text_fallback(p)


PARSERS: Dict[str, Callable[[Path], str]] = {
    ".txt": parse_txt_md,
    ".md": parse_txt_md,
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
    ".xlsx": parse_xlsx,
    ".csv": parse_csv,
    ".epub": parse_epub,
}


def parse_file(path: str) -> List[Dict[str, Any]]:
    p = Path(path)
    ext = p.suffix.lower()
    parser = PARSERS.get(ext, parse_txt_md)
    text = parser(p)
    chunks = _split_chunks(text)
    return [
        {"text": c, "meta": {"path": str(p), "ext": ext}} for c in chunks if c.strip()
    ]
