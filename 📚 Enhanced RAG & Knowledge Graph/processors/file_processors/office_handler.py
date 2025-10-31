"""
Office parsers for docx/xlsx/pptx for the emoji-named package.
Optional deps: python-docx, openpyxl, python-pptx.
"""

from __future__ import annotations

from pathlib import Path
from typing import Dict, List


def _md5_file(path: Path) -> str:
    import hashlib

    h = hashlib.md5()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()


def parse_docx(path: str | Path) -> List[Dict]:
    try:
        from docx import Document  # python-docx
    except Exception as e:
        raise ImportError(f"python-docx not available: {e}")
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(str(p))
    doc = Document(str(p))
    checksum = _md5_file(p)
    chunks: List[Dict] = []
    for i, para in enumerate(doc.paragraphs, start=1):
        text = (para.text or "").strip()
        if not text:
            continue
        chunks.append(
            {
                "text": text,
                "source": str(p),
                "page": i,
                "checksum": checksum,
            }
        )
    return chunks


def parse_xlsx(path: str | Path) -> List[Dict]:
    try:
        from openpyxl import load_workbook
    except Exception as e:
        raise ImportError(f"openpyxl not available: {e}")
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(str(p))
    wb = load_workbook(str(p), data_only=True, read_only=True)
    checksum = _md5_file(p)
    chunks: List[Dict] = []
    for ws in wb.worksheets:
        for row_idx, row in enumerate(ws.iter_rows(values_only=True), start=1):
            vals = [str(v) for v in row if v is not None]
            if not vals:
                continue
            text = " | ".join(vals)
            chunks.append(
                {
                    "text": text,
                    "source": f"{p}::{ws.title}",
                    "page": row_idx,
                    "checksum": checksum,
                }
            )
    return chunks


def parse_pptx(path: str | Path) -> List[Dict]:
    try:
        from pptx import Presentation  # python-pptx
    except Exception as e:
        raise ImportError(f"python-pptx not available: {e}")
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(str(p))
    prs = Presentation(str(p))
    checksum = _md5_file(p)
    chunks: List[Dict] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
            elif hasattr(shape, "text_frame") and shape.text_frame:
                texts.append(shape.text_frame.text or "")
        text = "\n".join(t.strip() for t in texts if t and t.strip())
        if not text:
            continue
        chunks.append(
            {
                "text": text,
                "source": str(p),
                "page": i,
                "checksum": checksum,
            }
        )
    return chunks
