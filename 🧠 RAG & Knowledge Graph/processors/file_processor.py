"""
文件处理器 - 支持全格式文件处理

支持的文件格式：
1. 办公文件：PDF, DOCX, XLSX, PPTX
2. 文本文件：TXT, MD, JSON, XML, CSV
3. 代码文件：PY, JS, JAVA, CPP, GO, RS等
4. 图片文件：JPG, PNG, GIF + OCR
5. 音频文件：MP3, WAV + 转文本
6. 视频文件：MP4, AVI + 转文本
"""

import os
import mimetypes
from pathlib import Path
from typing import Dict, Any, Optional, List
from datetime import datetime
import logging

# 设置日志
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class FileProcessor:
    """
    文件处理器
    
    负责：
    1. 识别文件类型
    2. 提取文件内容
    3. 处理元数据
    4. 支持批量处理
    """
    
    def __init__(self, config: Optional[Dict] = None):
        """初始化文件处理器"""
        self.config = config or self._get_default_config()
        self.supported_formats = self._get_supported_formats()
        
        logger.info("📁 文件处理器初始化完成")
        logger.info(f"   支持格式: {len(self.supported_formats)}种")
    
    def _get_default_config(self) -> Dict:
        """获取默认配置"""
        return {
            "max_file_size": 100 * 1024 * 1024,  # 100MB
            "ocr_enabled": True,
            "ocr_language": "chi_sim+eng",
            "audio_to_text": True,
            "video_to_text": False  # 视频转文本较慢，默认关闭
        }
    
    def _get_supported_formats(self) -> Dict[str, List[str]]:
        """获取支持的文件格式"""
        return {
            "documents": [".pdf", ".docx", ".xlsx", ".pptx", ".doc", ".xls", ".ppt"],
            "text": [".txt", ".md", ".markdown", ".rst"],
            "code": [".py", ".js", ".java", ".cpp", ".c", ".h", ".go", ".rs", ".ts", 
                    ".jsx", ".tsx", ".vue", ".html", ".css", ".sql", ".sh"],
            "data": [".json", ".xml", ".csv", ".yaml", ".yml", ".toml"],
            "images": [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp"],
            "audio": [".mp3", ".wav", ".m4a", ".flac", ".ogg"],
            "video": [".mp4", ".avi", ".mkv", ".mov", ".wmv"]
        }
    
    def process_file(self, file_path: str) -> Dict[str, Any]:
        """
        处理单个文件
        
        Args:
            file_path: 文件路径
            
        Returns:
            处理结果，包含内容和元数据
        """
        logger.info(f"\n📄 处理文件: {file_path}")
        
        # 检查文件是否存在
        if not os.path.exists(file_path):
            return self._error_result(f"文件不存在: {file_path}")
        
        # 检查文件大小
        file_size = os.path.getsize(file_path)
        if file_size > self.config["max_file_size"]:
            return self._error_result(f"文件过大: {file_size / 1024 / 1024:.2f}MB")
        
        # 识别文件类型
        file_type, file_category = self._identify_file_type(file_path)
        logger.info(f"   类型: {file_type} | 类别: {file_category}")
        
        # 提取内容
        try:
            content = self._extract_content(file_path, file_category)
            
            # 提取元数据
            metadata = self._extract_metadata(file_path, file_type, file_category)
            
            result = {
                "success": True,
                "file_path": file_path,
                "file_name": os.path.basename(file_path),
                "file_type": file_type,
                "file_category": file_category,
                "file_size": file_size,
                "content": content,
                "content_length": len(content),
                "metadata": metadata,
                "processed_at": datetime.now().isoformat()
            }
            
            logger.info(f"   ✅ 处理成功，内容长度: {len(content)}字符")
            return result
            
        except Exception as e:
            logger.error(f"   ❌ 处理失败: {str(e)}")
            return self._error_result(str(e))
    
    def _identify_file_type(self, file_path: str) -> tuple:
        """识别文件类型"""
        ext = Path(file_path).suffix.lower()
        
        # 遍历支持的格式
        for category, extensions in self.supported_formats.items():
            if ext in extensions:
                return ext, category
        
        # 如果扩展名不在支持列表中，尝试MIME类型
        mime_type, _ = mimetypes.guess_type(file_path)
        if mime_type:
            if mime_type.startswith("text/"):
                return ext, "text"
            elif mime_type.startswith("image/"):
                return ext, "images"
            elif mime_type.startswith("audio/"):
                return ext, "audio"
            elif mime_type.startswith("video/"):
                return ext, "video"
        
        return ext, "unknown"
    
    def _extract_content(self, file_path: str, file_category: str) -> str:
        """
        根据文件类别提取内容
        
        Args:
            file_path: 文件路径
            file_category: 文件类别
            
        Returns:
            提取的文本内容
        """
        if file_category == "documents":
            return self._extract_document_content(file_path)
        elif file_category == "text":
            return self._extract_text_content(file_path)
        elif file_category == "code":
            return self._extract_code_content(file_path)
        elif file_category == "data":
            return self._extract_data_content(file_path)
        elif file_category == "images":
            return self._extract_image_content(file_path)
        elif file_category == "audio":
            return self._extract_audio_content(file_path)
        elif file_category == "video":
            return self._extract_video_content(file_path)
        else:
            return self._extract_raw_content(file_path)
    
    def _extract_document_content(self, file_path: str) -> str:
        """提取办公文档内容"""
        ext = Path(file_path).suffix.lower()
        
        if ext == ".pdf":
            return self._extract_pdf_content(file_path)
        elif ext in [".docx", ".doc"]:
            return self._extract_word_content(file_path)
        elif ext in [".xlsx", ".xls"]:
            return self._extract_excel_content(file_path)
        elif ext in [".pptx", ".ppt"]:
            return self._extract_ppt_content(file_path)
        else:
            return ""
    
    def _extract_pdf_content(self, file_path: str) -> str:
        """提取PDF内容"""
        try:
            # 尝试使用PyPDF2
            import PyPDF2
            content = []
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        content.append(text)
            
            result = "\n".join(content)
            
            # 如果提取内容很少，可能是图片PDF，尝试OCR
            if len(result) < 100 and self.config.get("ocr_enabled"):
                logger.info("   PDF文本较少，尝试OCR...")
                # TODO: 实现OCR
                pass
            
            return result
            
        except ImportError:
            logger.warning("   PyPDF2未安装，无法处理PDF")
            return ""
        except Exception as e:
            logger.error(f"   PDF处理错误: {e}")
            return ""
    
    def _extract_word_content(self, file_path: str) -> str:
        """提取Word文档内容"""
        try:
            from docx import Document
            doc = Document(file_path)
            content = []
            
            # 提取段落
            for para in doc.paragraphs:
                if para.text.strip():
                    content.append(para.text)
            
            # 提取表格
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    if row_text.strip():
                        content.append(row_text)
            
            return "\n".join(content)
            
        except ImportError:
            logger.warning("   python-docx未安装，无法处理Word")
            return ""
        except Exception as e:
            logger.error(f"   Word处理错误: {e}")
            return ""
    
    def _extract_excel_content(self, file_path: str) -> str:
        """提取Excel内容"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True)
            content = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                content.append(f"\n=== {sheet_name} ===\n")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        content.append(row_text)
            
            return "\n".join(content)
            
        except ImportError:
            logger.warning("   openpyxl未安装，无法处理Excel")
            return ""
        except Exception as e:
            logger.error(f"   Excel处理错误: {e}")
            return ""
    
    def _extract_ppt_content(self, file_path: str) -> str:
        """提取PowerPoint内容"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            content = []
            
            for i, slide in enumerate(prs.slides, 1):
                content.append(f"\n=== Slide {i} ===\n")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content.append(shape.text)
            
            return "\n".join(content)
            
        except ImportError:
            logger.warning("   python-pptx未安装，无法处理PowerPoint")
            return ""
        except Exception as e:
            logger.error(f"   PowerPoint处理错误: {e}")
            return ""
    
    def _extract_text_content(self, file_path: str) -> str:
        """提取文本文件内容"""
        try:
            # 尝试多种编码
            encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            
            # 如果所有编码都失败，使用二进制模式
            with open(file_path, 'rb') as f:
                return f.read().decode('utf-8', errors='ignore')
                
        except Exception as e:
            logger.error(f"   文本文件读取错误: {e}")
            return ""
    
    def _extract_code_content(self, file_path: str) -> str:
        """提取代码文件内容"""
        # 代码文件也是文本文件，直接调用文本提取
        content = self._extract_text_content(file_path)
        
        # 添加代码文件的特殊标记
        ext = Path(file_path).suffix.lower()
        language = self._get_language_from_extension(ext)
        
        return f"```{language}\n{content}\n```"
    
    def _get_language_from_extension(self, ext: str) -> str:
        """根据扩展名获取编程语言"""
        language_map = {
            ".py": "python",
            ".js": "javascript",
            ".ts": "typescript",
            ".java": "java",
            ".cpp": "cpp",
            ".c": "c",
            ".h": "c",
            ".go": "go",
            ".rs": "rust",
            ".sql": "sql",
            ".sh": "bash",
            ".html": "html",
            ".css": "css",
            ".vue": "vue",
            ".jsx": "jsx",
            ".tsx": "tsx"
        }
        return language_map.get(ext, "text")
    
    def _extract_data_content(self, file_path: str) -> str:
        """提取数据文件内容"""
        ext = Path(file_path).suffix.lower()
        
        if ext == ".json":
            import json
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    return json.dumps(data, ensure_ascii=False, indent=2)
            except Exception as e:
                logger.error(f"   JSON解析错误: {e}")
                return self._extract_text_content(file_path)
        
        elif ext == ".xml":
            return self._extract_text_content(file_path)
        
        elif ext == ".csv":
            try:
                import pandas as pd
                df = pd.read_csv(file_path)
                return df.to_string()
            except ImportError:
                return self._extract_text_content(file_path)
            except Exception as e:
                logger.error(f"   CSV解析错误: {e}")
                return self._extract_text_content(file_path)
        
        elif ext in [".yaml", ".yml"]:
            import yaml
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    data = yaml.safe_load(f)
                    return yaml.dump(data, allow_unicode=True)
            except Exception as e:
                logger.error(f"   YAML解析错误: {e}")
                return self._extract_text_content(file_path)
        
        else:
            return self._extract_text_content(file_path)
    
    def _extract_image_content(self, file_path: str) -> str:
        """提取图片内容（OCR）"""
        if not self.config.get("ocr_enabled"):
            return f"[图片文件: {os.path.basename(file_path)}]"
        
        try:
            # 尝试使用pytesseract进行OCR
            from PIL import Image
            import pytesseract
            
            image = Image.open(file_path)
            text = pytesseract.image_to_string(
                image, 
                lang=self.config.get("ocr_language", "chi_sim+eng")
            )
            
            if text.strip():
                return f"[图片OCR结果]\n{text}"
            else:
                return f"[图片文件: {os.path.basename(file_path)}, OCR未识别到文本]"
                
        except ImportError:
            logger.warning("   pytesseract未安装，无法进行OCR")
            return f"[图片文件: {os.path.basename(file_path)}]"
        except Exception as e:
            logger.error(f"   OCR错误: {e}")
            return f"[图片文件: {os.path.basename(file_path)}]"
    
    def _extract_audio_content(self, file_path: str) -> str:
        """提取音频内容（转文本）"""
        if not self.config.get("audio_to_text"):
            return f"[音频文件: {os.path.basename(file_path)}]"
        
        try:
            # TODO: 实现音频转文本（需要语音识别库）
            # 可以使用：
            # - SpeechRecognition
            # - Whisper
            # - 云服务API
            
            logger.info("   音频转文本功能待实现")
            return f"[音频文件: {os.path.basename(file_path)}]"
            
        except Exception as e:
            logger.error(f"   音频处理错误: {e}")
            return f"[音频文件: {os.path.basename(file_path)}]"
    
    def _extract_video_content(self, file_path: str) -> str:
        """提取视频内容（转文本）"""
        if not self.config.get("video_to_text"):
            return f"[视频文件: {os.path.basename(file_path)}]"
        
        try:
            # TODO: 实现视频转文本
            # 1. 提取音频轨道
            # 2. 音频转文本
            # 3. 可选：关键帧OCR
            
            logger.info("   视频转文本功能待实现")
            return f"[视频文件: {os.path.basename(file_path)}]"
            
        except Exception as e:
            logger.error(f"   视频处理错误: {e}")
            return f"[视频文件: {os.path.basename(file_path)}]"
    
    def _extract_raw_content(self, file_path: str) -> str:
        """提取未知格式文件的原始内容"""
        try:
            with open(file_path, 'rb') as f:
                content = f.read(1024)  # 只读前1KB
                return content.decode('utf-8', errors='ignore')
        except Exception as e:
            logger.error(f"   原始内容提取错误: {e}")
            return f"[无法读取文件: {os.path.basename(file_path)}]"
    
    def _extract_metadata(self, file_path: str, file_type: str, file_category: str) -> Dict:
        """提取文件元数据"""
        stat = os.stat(file_path)
        
        return {
            "file_name": os.path.basename(file_path),
            "file_path": file_path,
            "file_type": file_type,
            "file_category": file_category,
            "file_size": stat.st_size,
            "created_at": datetime.fromtimestamp(stat.st_ctime).isoformat(),
            "modified_at": datetime.fromtimestamp(stat.st_mtime).isoformat(),
            "mime_type": mimetypes.guess_type(file_path)[0]
        }
    
    def _error_result(self, error_message: str) -> Dict:
        """创建错误结果"""
        return {
            "success": False,
            "error": error_message,
            "processed_at": datetime.now().isoformat()
        }
    
    def process_directory(self, directory_path: str, recursive: bool = True) -> List[Dict]:
        """
        批量处理目录中的文件
        
        Args:
            directory_path: 目录路径
            recursive: 是否递归处理子目录
            
        Returns:
            所有文件的处理结果列表
        """
        logger.info(f"\n📁 批量处理目录: {directory_path}")
        
        results = []
        
        if recursive:
            for root, dirs, files in os.walk(directory_path):
                for file in files:
                    file_path = os.path.join(root, file)
                    result = self.process_file(file_path)
                    results.append(result)
        else:
            for file in os.listdir(directory_path):
                file_path = os.path.join(directory_path, file)
                if os.path.isfile(file_path):
                    result = self.process_file(file_path)
                    results.append(result)
        
        success_count = sum(1 for r in results if r.get("success"))
        logger.info(f"\n✅ 批量处理完成: {success_count}/{len(results)} 成功")
        
        return results
    
    def get_supported_formats_info(self) -> Dict:
        """获取支持的格式信息"""
        total = sum(len(exts) for exts in self.supported_formats.values())
        
        return {
            "total_formats": total,
            "categories": {
                category: {
                    "count": len(exts),
                    "extensions": exts
                }
                for category, exts in self.supported_formats.items()
            }
        }


def test_file_processor():
    """测试文件处理器"""
    print("="*70)
    print("  文件处理器测试")
    print("="*70)
    
    processor = FileProcessor()
    
    # 显示支持的格式
    formats_info = processor.get_supported_formats_info()
    print(f"\n支持的格式数量: {formats_info['total_formats']}")
    for category, info in formats_info['categories'].items():
        print(f"  {category}: {info['count']}种 - {', '.join(info['extensions'][:5])}...")
    
    # 测试文本文件处理
    test_file = "test.txt"
    if os.path.exists(test_file):
        result = processor.process_file(test_file)
        print(f"\n测试结果:")
        print(f"  成功: {result.get('success')}")
        print(f"  内容长度: {result.get('content_length')}字符")
    
    print("\n✅ 文件处理器测试完成！")


if __name__ == "__main__":
    test_file_processor()






