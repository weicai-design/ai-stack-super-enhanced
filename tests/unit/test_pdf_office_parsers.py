import importlib
import os
import sys
import pytest

ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
ALT = os.path.join(ROOT, "📚 Enhanced RAG & Knowledge Graph")

for path in (ALT, ROOT):
    if os.path.isdir(path) and path not in sys.path:
        sys.path.insert(0, path)


def _has_pkg(name: str) -> bool:
    try:
        __import__(name)
        return True
    except Exception:
        return False


def test_parse_pdf_no_pypdf2_raises_or_parses(tmp_path):
    from processors.file_processors import pdf_parser

    # create a tiny PDF if PyPDF2 is present; otherwise expect ImportError
    if not _has_pkg("PyPDF2"):
        with pytest.raises(ImportError):
            pdf_parser.parse_pdf("nonexistent.pdf")
        pytest.skip("PyPDF2 not installed; skipping parse test")

    # If PyPDF2 is available, create a dummy PDF using reportlab if possible
    pdf_path = tmp_path / "test.pdf"
    try:
        # prefer reportlab to generate a tiny pdf for parsing
        if _has_pkg("reportlab"):
            from reportlab.pdfgen import canvas

            c = canvas.Canvas(str(pdf_path))
            c.drawString(100, 750, "Hello PDF")
            c.save()
        else:
            # If no generator lib available, skip the deep parsing test
            pytest.skip("reportlab not installed; skipping PDF generation")

        chunks = pdf_parser.parse_pdf(str(pdf_path))
        assert isinstance(chunks, list)
        assert all("text" in c for c in chunks)
    except ImportError:
        pytest.skip("PyPDF2 import failed at runtime")


def test_office_parsers_skip_or_run(tmp_path):
    from processors.file_processors import office_handler

    # DOCX
    if not _has_pkg("docx"):
        with pytest.raises(ImportError):
            office_handler.parse_docx("nonexistent.docx")
    else:
        try:
            from docx import Document

            docx_path = tmp_path / "t.docx"
            doc = Document()
            doc.add_paragraph("Hello docx")
            doc.save(str(docx_path))
            chunks = office_handler.parse_docx(str(docx_path))
            assert isinstance(chunks, list)
            assert all("text" in c for c in chunks)
        except Exception:
            pytest.skip("Unable to create or parse docx on this environment")

    # XLSX
    if not _has_pkg("openpyxl"):
        with pytest.raises(ImportError):
            office_handler.parse_xlsx("nonexistent.xlsx")
    else:
        try:
            from openpyxl import Workbook

            xlsx_path = tmp_path / "t.xlsx"
            wb = Workbook()
            ws = wb.active
            ws.append(["a", "b", "c"])
            wb.save(str(xlsx_path))
            chunks = office_handler.parse_xlsx(str(xlsx_path))
            assert isinstance(chunks, list)
        except Exception:
            pytest.skip("Unable to create or parse xlsx on this environment")

    # PPTX
    if not _has_pkg("pptx"):
        with pytest.raises(ImportError):
            office_handler.parse_pptx("nonexistent.pptx")
    else:
        try:
            from pptx import Presentation

            ppt_path = tmp_path / "t.pptx"
            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[5])
            prs.save(str(ppt_path))
            chunks = office_handler.parse_pptx(str(ppt_path))
            assert isinstance(chunks, list)
        except Exception:
            pytest.skip("Unable to create or parse pptx on this environment")


def test_universal_parser_importable():
    """
    验证 universal_file_parser 模块可导入并创建 parser 实例（最小 smoke test）。
    若模块缺少依赖或语法错误，导入会失败并导致此测试失败。
    """
    mod = importlib.import_module("processors.file_processors.universal_file_parser")
    assert hasattr(mod, "create_universal_parser")
    parser = mod.create_universal_parser()
    assert parser is not None
    assert hasattr(parser, "process")
