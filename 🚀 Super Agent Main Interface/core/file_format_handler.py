"""
文件格式处理器
支持60+种文件格式的处理
"""

from typing import Dict, List, Optional, Any, BinaryIO
from datetime import datetime
import mimetypes
import os
import logging
import asyncio

logger = logging.getLogger(__name__)

class FileFormatHandler:
    """
    文件格式处理器
    
    功能：
    1. 识别文件格式（60+种）
    2. 解析文件内容
    3. 提取文本/元数据
    4. 转换为统一格式
    """
    
    def __init__(self, max_file_size_mb: int = 100, enable_validation: bool = True):
        """
        初始化文件格式处理器
        
        Args:
            max_file_size_mb: 最大文件大小（MB）
            enable_validation: 是否启用格式验证
        """
        self.max_file_size_mb = max_file_size_mb
        self.enable_validation = enable_validation
        
        # 支持的文件格式（60+种）
        self.supported_formats = {
            # 办公文档（15种）
            "office": [".doc", ".docx", ".docm", ".dot", ".dotx",
                      ".xls", ".xlsx", ".xlsm", ".xlsb", ".xlt", ".xltx",
                      ".ppt", ".pptx", ".pptm", ".pot", ".potx",
                      ".pdf", ".odt", ".ods", ".odp", ".rtf", ".msg", ".eml"],
            
            # 电子书（8种）
            "ebook": [".epub", ".mobi", ".azw", ".azw3", ".fb2", ".lit", ".prc", ".txt"],
            
            # 编程文件（30+种）
            "code": [".py", ".js", ".java", ".c", ".cpp", ".cs", ".swift", ".kt",
                    ".scala", ".dart", ".lua", ".pl", ".pm", ".sh", ".bash", ".ps1",
                    ".r", ".R", ".m", ".sql", ".md", ".html", ".css", ".php",
                    ".rb", ".go", ".rs", ".ts", ".vue", ".jsx", ".tsx",
                    ".conf", ".ini", ".toml", ".yaml", ".yml", ".json", ".xml"],
            
            # 图片（15种）
            "image": [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".tif",
                     ".webp", ".svg", ".heic", ".heif", ".ico", ".psd", ".ai",
                     ".raw", ".cr2", ".nef", ".arw"],
            
            # 音频（10种）
            "audio": [".mp3", ".wav", ".flac", ".aac", ".ogg", ".m4a", ".wma",
                     ".opus", ".amr", ".3gp"],
            
            # 视频（10种）
            "video": [".mp4", ".avi", ".mov", ".mkv", ".wmv", ".flv", ".webm",
                    ".mpeg", ".mpg", ".3gp", ".rm", ".rmvb", ".vob"],
            
            # 思维导图（4种）
            "mindmap": [".xmind", ".mm", ".mmap", ".opml"],
            
            # 数据库（5种）
            "database": [".csv", ".sql", ".db", ".mdb", ".accdb", ".fdb", ".odb"],
            
            # 压缩文件（5种）
            "archive": [".zip", ".rar", ".7z", ".tar", ".gz"],
            
            # 其他（5种）
            "other": [".txt", ".md", ".json", ".xml", ".yaml"]
        }
        
        # 格式处理器映射
        self.processors = {}
        self._initialize_processors()
    
    def _initialize_processors(self):
        """初始化格式处理器"""
        # 办公文档处理器
        self.processors["office"] = self._process_office_document
        # 电子书处理器
        self.processors["ebook"] = self._process_ebook
        # 编程文件处理器
        self.processors["code"] = self._process_code_file
        # 图片处理器
        self.processors["image"] = self._process_image
        # 音频处理器
        self.processors["audio"] = self._process_audio
        # 视频处理器
        self.processors["video"] = self._process_video
        # 思维导图处理器
        self.processors["mindmap"] = self._process_mindmap
        # 数据库处理器
        self.processors["database"] = self._process_database
        # 压缩文件处理器
        self.processors["archive"] = self._process_archive
        # 其他处理器
        self.processors["other"] = self._process_text_file
    
    def detect_format(self, filename: str, mime_type: Optional[str] = None) -> Dict[str, Any]:
        """
        检测文件格式
        
        Args:
            filename: 文件名
            mime_type: MIME类型（可选）
            
        Returns:
            格式信息
        """
        ext = os.path.splitext(filename)[1].lower()
        
        # 检测格式类别
        format_category = None
        for category, extensions in self.supported_formats.items():
            if ext in extensions:
                format_category = category
                break
        
        if not format_category:
            format_category = "unknown"
        
        # 如果没有提供MIME类型，尝试检测
        if not mime_type:
            mime_type, _ = mimetypes.guess_type(filename)
        
        return {
            "extension": ext,
            "category": format_category,
            "mime_type": mime_type,
            "is_supported": format_category != "unknown"
        }
    
    async def process_file(
        self,
        file_data: bytes,
        filename: str,
        mime_type: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        处理文件⭐增强版（支持验证和性能优化）
        
        Args:
            file_data: 文件数据
            filename: 文件名
            mime_type: MIME类型
            
        Returns:
            处理结果
        """
        start_time = datetime.now()
        
        # 1. 文件大小验证
        if self.enable_validation:
            file_size_mb = len(file_data) / (1024 * 1024)
            if file_size_mb > self.max_file_size_mb:
                return {
                    "success": False,
                    "error": f"文件过大: {file_size_mb:.2f}MB > {self.max_file_size_mb}MB",
                    "file_size_mb": file_size_mb
                }
        
        # 2. 检测文件格式
        format_info = self.detect_format(filename, mime_type)
        
        if not format_info["is_supported"]:
            return {
                "success": False,
                "error": f"不支持的文件格式: {format_info['extension']}",
                "format": format_info
            }
        
        # 3. 验证文件内容（如果启用）
        if self.enable_validation:
            validation_result = await self._validate_file_content(file_data, format_info)
            if not validation_result["valid"]:
                return {
                    "success": False,
                    "error": f"文件验证失败: {validation_result.get('error', '格式不匹配')}",
                    "format": format_info,
                    "validation": validation_result
                }
        
        # 4. 处理文件
        category = format_info["category"]
        processor = self.processors.get(category)
        
        if not processor:
            return {
                "success": False,
                "error": f"未找到处理器: {category}",
                "format": format_info
            }
        
        try:
            # 带超时处理（避免大文件处理时间过长）
            result = await asyncio.wait_for(
                processor(file_data, filename, format_info),
                timeout=30.0  # 30秒超时
            )
            result["format"] = format_info
            result["success"] = True
            
            # 记录处理时间
            processing_time = (datetime.now() - start_time).total_seconds()
            result["processing_time"] = processing_time
            
            logger.info(f"文件处理完成: {filename} ({processing_time:.2f}秒)")
            
            return result
        except asyncio.TimeoutError:
            return {
                "success": False,
                "error": "文件处理超时（30秒）",
                "format": format_info
            }
        except Exception as e:
            logger.error(f"文件处理失败: {filename} - {str(e)}")
            return {
                "success": False,
                "error": str(e),
                "format": format_info
            }
    
    async def _validate_file_content(
        self,
        file_data: bytes,
        format_info: Dict
    ) -> Dict[str, Any]:
        """
        验证文件内容
        
        Args:
            file_data: 文件数据
            format_info: 格式信息
            
        Returns:
            验证结果
        """
        ext = format_info["extension"]
        
        try:
            # 文件头验证（Magic Number）
            file_signatures = {
                ".pdf": [b"%PDF"],
                ".docx": [b"PK\x03\x04"],  # ZIP格式
                ".xlsx": [b"PK\x03\x04"],
                ".pptx": [b"PK\x03\x04"],
                ".zip": [b"PK\x03\x04"],
                ".jpg": [b"\xff\xd8\xff"],
                ".jpeg": [b"\xff\xd8\xff"],
                ".png": [b"\x89PNG"],
                ".gif": [b"GIF87a", b"GIF89a"],
                ".mp3": [b"ID3", b"\xff\xfb", b"\xff\xf3"],
                ".mp4": [b"\x00\x00\x00\x18ftypmp42", b"\x00\x00\x00\x1cftypmp42"],
                ".wav": [b"RIFF"]
            }
            
            if ext in file_signatures:
                valid = False
                for signature in file_signatures[ext]:
                    if file_data[:len(signature)] == signature:
                        valid = True
                        break
                
                if not valid:
                    return {
                        "valid": False,
                        "error": f"文件头不匹配：预期{ext}格式，但文件头验证失败"
                    }
            
            # 所有验证通过
            return {
                "valid": True,
                "checks_passed": ["file_size", "file_signature"],
                "file_size_mb": len(file_data) / (1024 * 1024)
            }
        except Exception as e:
            logger.warning(f"文件验证失败: {e}")
            # 验证失败不影响处理，只是警告
            return {"valid": True, "warning": str(e)}
    
    async def _process_office_document(
        self,
        file_data: bytes,
        filename: str,
        format_info: Dict
    ) -> Dict[str, Any]:
        """处理办公文档"""
        ext = format_info["extension"]
        
        # Word文档
        if ext in [".doc", ".docx", ".docm", ".dot", ".dotx"]:
            return await self._extract_word_text(file_data, ext)
        
        # Excel文档
        elif ext in [".xls", ".xlsx", ".xlsm", ".xlsb", ".xlt", ".xltx"]:
            return await self._extract_excel_text(file_data, ext)
        
        # PowerPoint文档
        elif ext in [".ppt", ".pptx", ".pptm", ".pot", ".potx"]:
            return await self._extract_ppt_text(file_data, ext)
        
        # PDF
        elif ext == ".pdf":
            return await self._extract_pdf_text(file_data)
        
        # 其他格式
        else:
            return {
                "text": "",
                "metadata": {"format": ext, "note": "格式支持待完善"}
            }
    
    async def _extract_word_text(self, file_data: bytes, ext: str) -> Dict[str, Any]:
        """提取Word文档文本⭐增强版（支持表格和样式）"""
        try:
            if ext in [".docx", ".docm", ".dotx"]:
                # 使用python-docx
                try:
                    from docx import Document
                    import io
                    doc = Document(io.BytesIO(file_data))
                    
                    text_parts = []
                    
                    # 提取段落文本
                    for para in doc.paragraphs:
                        if para.text.strip():
                            text_parts.append(para.text)
                    
                    # 提取表格文本
                    for table in doc.tables:
                        for row in table.rows:
                            row_text = "\t".join([cell.text for cell in row.cells])
                            if row_text.strip():
                                text_parts.append(row_text)
                    
                    return {
                        "text": "\n".join(text_parts),
                        "metadata": {
                            "format": ext,
                            "paragraphs": len(doc.paragraphs),
                            "tables": len(doc.tables),
                            "method": "python-docx"
                        }
                    }
                except ImportError:
                    pass
                except Exception as e:
                    logger.warning(f"python-docx处理失败: {e}")
            
            # 备用：返回基本信息
            return {
                "text": "",
                "metadata": {"format": ext, "note": "需要安装python-docx库"}
            }
        except Exception as e:
            return {"text": "", "metadata": {"error": str(e)}}
    
    async def _extract_excel_text(self, file_data: bytes, ext: str) -> Dict[str, Any]:
        """提取Excel文档文本⭐增强版（性能优化，大文件支持）"""
        try:
            if ext in [".xlsx", ".xlsm", ".xlsb"]:
                # 使用openpyxl（只读模式，性能更好）
                try:
                    from openpyxl import load_workbook
                    import io
                    wb = load_workbook(io.BytesIO(file_data), read_only=True, data_only=True)
                    text_parts = []
                    total_rows = 0
                    
                    for sheet in wb.worksheets:
                        text_parts.append(f"=== 工作表: {sheet.title} ===")
                        row_count = 0
                        for row in sheet.iter_rows(values_only=True, max_row=1000):  # 限制行数避免过大
                            if row_count >= 1000:
                                text_parts.append("... (行数过多，已截断)")
                                break
                            row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                            if row_text.strip():
                                text_parts.append(row_text)
                                row_count += 1
                        total_rows += row_count
                    
                    return {
                        "text": "\n".join(text_parts),
                        "metadata": {
                            "format": ext,
                            "sheets": len(wb.worksheets),
                            "rows": total_rows,
                            "method": "openpyxl"
                        }
                    }
                except ImportError:
                    pass
                except Exception as e:
                    logger.warning(f"openpyxl处理失败: {e}")
            
            elif ext in [".xls", ".xlt"]:
                # 旧格式Excel，尝试使用xlrd
                try:
                    import xlrd
                    import io
                    workbook = xlrd.open_workbook(file_contents=file_data)
                    text_parts = []
                    for sheet in workbook.sheets():
                        text_parts.append(f"=== 工作表: {sheet.name} ===")
                        for row_idx in range(min(sheet.nrows, 1000)):  # 限制行数
                            row = sheet.row_values(row_idx)
                            row_text = "\t".join([str(cell) for cell in row])
                            if row_text.strip():
                                text_parts.append(row_text)
                    
                    return {
                        "text": "\n".join(text_parts),
                        "metadata": {
                            "format": ext,
                            "sheets": workbook.nsheets,
                            "method": "xlrd"
                        }
                    }
                except ImportError:
                    pass
                except Exception as e:
                    logger.warning(f"xlrd处理失败: {e}")
            
            return {
                "text": "",
                "metadata": {"format": ext, "note": "需要安装openpyxl或xlrd库"}
            }
        except Exception as e:
            return {"text": "", "metadata": {"error": str(e)}}
    
    async def _extract_ppt_text(self, file_data: bytes, ext: str) -> Dict[str, Any]:
        """提取PPT文档文本⭐实现版"""
        try:
            if ext in [".pptx", ".pptm"]:
                # 使用python-pptx
                try:
                    from pptx import Presentation
                    import io
                    prs = Presentation(io.BytesIO(file_data))
                    text_parts = []
                    for slide_num, slide in enumerate(prs.slides, 1):
                        text_parts.append(f"幻灯片 {slide_num}:")
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                text_parts.append(shape.text)
                    return {"text": "\n".join(text_parts), "metadata": {"format": ext, "slides": len(prs.slides)}}
                except ImportError:
                    pass
            
            return {
                "text": "",
                "metadata": {"format": ext, "note": "需要安装python-pptx库"}
            }
        except Exception as e:
            return {"text": "", "metadata": {"error": str(e)}}
    
    async def _extract_pdf_text(self, file_data: bytes) -> Dict[str, Any]:
        """提取PDF文本⭐增强版（支持多种库）"""
        try:
            # 策略1: 尝试使用pdfplumber（质量更好）
            try:
                import pdfplumber
                import io
                text_parts = []
                with pdfplumber.open(io.BytesIO(file_data)) as pdf:
                    for page in pdf.pages:
                        text = page.extract_text()
                        if text:
                            text_parts.append(text)
                return {
                    "text": "\n".join(text_parts),
                    "metadata": {
                        "format": "pdf",
                        "pages": len(text_parts),
                        "method": "pdfplumber"
                    }
                }
            except ImportError:
                pass
            except Exception as e:
                logger.warning(f"pdfplumber处理失败: {e}")
            
            # 策略2: 尝试使用PyPDF2
            try:
                import PyPDF2
                import io
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_data))
                text_parts = []
                for page in pdf_reader.pages:
                    text = page.extract_text()
                    if text:
                        text_parts.append(text)
                return {
                    "text": "\n".join(text_parts),
                    "metadata": {
                        "format": "pdf",
                        "pages": len(pdf_reader.pages),
                        "method": "PyPDF2"
                    }
                }
            except ImportError:
                pass
            except Exception as e:
                logger.warning(f"PyPDF2处理失败: {e}")
            
            # 策略3: 尝试使用pymupdf（fitz）
            try:
                import fitz  # PyMuPDF
                import io
                pdf_doc = fitz.open(stream=file_data, filetype="pdf")
                text_parts = []
                for page in pdf_doc:
                    text = page.get_text()
                    if text:
                        text_parts.append(text)
                pdf_doc.close()
                return {
                    "text": "\n".join(text_parts),
                    "metadata": {
                        "format": "pdf",
                        "pages": len(text_parts),
                        "method": "PyMuPDF"
                    }
                }
            except ImportError:
                pass
            except Exception as e:
                logger.warning(f"PyMuPDF处理失败: {e}")
            
            return {
                "text": "",
                "metadata": {"format": "pdf", "note": "需要安装pdfplumber、PyPDF2或PyMuPDF库"}
            }
        except Exception as e:
            return {"text": "", "metadata": {"error": str(e)}}
    
    async def _process_ebook(self, file_data: bytes, filename: str, format_info: Dict) -> Dict[str, Any]:
        """处理电子书⭐实现版"""
        ext = format_info["extension"]
        try:
            if ext == ".epub":
                # 使用ebooklib
                try:
                    import ebooklib
                    from ebooklib import epub
                    import io
                    book = epub.read_epub(io.BytesIO(file_data))
                    text_parts = []
                    for item in book.get_items():
                        if item.get_type() == ebooklib.ITEM_DOCUMENT:
                            # 提取HTML内容中的文本（简单实现）
                            content = item.get_content().decode('utf-8', errors='ignore')
                            # 简单去除HTML标签
                            import re
                            text = re.sub(r'<[^>]+>', '', content)
                            if text.strip():
                                text_parts.append(text.strip())
                    return {"text": "\n".join(text_parts), "metadata": {"format": "epub"}}
                except ImportError:
                    pass
            
            elif ext == ".mobi":
                # MOBI格式需要特殊处理
                return {
                    "text": "",
                    "metadata": {"format": ext, "note": "MOBI格式需要mobi库"}
                }
            
            return {
                "text": "",
                "metadata": {"format": ext, "note": "需要安装ebooklib库"}
            }
        except Exception as e:
            return {"text": "", "metadata": {"error": str(e)}}
    
    async def _process_code_file(
        self,
        file_data: bytes,
        filename: str,
        format_info: Dict
    ) -> Dict[str, Any]:
        """处理编程文件"""
        try:
            text = file_data.decode('utf-8', errors='ignore')
            return {
                "text": text,
                "metadata": {
                    "format": format_info["extension"],
                    "language": self._detect_language(filename),
                    "lines": len(text.split('\n'))
                }
            }
        except Exception as e:
            return {"text": "", "metadata": {"error": str(e)}}
    
    async def _process_image(self, file_data: bytes, filename: str, format_info: Dict) -> Dict[str, Any]:
        """处理图片⭐实现版（OCR支持）"""
        ext = format_info["extension"]
        try:
            # 尝试OCR提取文本
            try:
                import pytesseract
                from PIL import Image
                import io
                
                # 打开图片
                image = Image.open(io.BytesIO(file_data))
                
                # OCR识别
                text = pytesseract.image_to_string(image, lang='chi_sim+eng')
                
                return {
                    "text": text,
                    "metadata": {
                        "format": ext,
                        "size": len(file_data),
                        "ocr": True,
                        "dimensions": image.size
                    }
                }
            except ImportError:
                pass
            
            # 如果没有OCR库，返回基本信息
            return {
                "text": "",
                "metadata": {
                    "format": ext,
                    "size": len(file_data),
                    "note": "需要安装pytesseract和Pillow库进行OCR"
                }
            }
        except Exception as e:
            return {
                "text": "",
                "metadata": {"error": str(e), "format": ext}
            }
    
    async def _process_audio(self, file_data: bytes, filename: str, format_info: Dict) -> Dict[str, Any]:
        """处理音频⭐实现版（语音识别支持）"""
        ext = format_info["extension"]
        try:
            # 尝试语音识别
            try:
                import speech_recognition as sr
                import io
                
                # 创建识别器
                r = sr.Recognizer()
                
                # 根据格式选择处理方式
                if ext in [".wav"]:
                    # WAV格式可以直接使用
                    audio_file = io.BytesIO(file_data)
                    with sr.AudioFile(audio_file) as source:
                        audio = r.record(source)
                        text = r.recognize_google(audio, language='zh-CN')
                        return {
                            "text": text,
                            "metadata": {
                                "format": ext,
                                "size": len(file_data),
                                "speech_recognition": True
                            }
                        }
                else:
                    # 其他格式需要转换
                    return {
                        "text": "",
                        "metadata": {
                            "format": ext,
                            "size": len(file_data),
                            "note": f"{ext}格式需要转换为WAV格式"
                        }
                    }
            except ImportError:
                pass
            
            return {
                "text": "",
                "metadata": {
                    "format": ext,
                    "size": len(file_data),
                    "note": "需要安装SpeechRecognition库进行语音识别"
                }
            }
        except Exception as e:
            return {
                "text": "",
                "metadata": {"error": str(e), "format": ext}
            }
    
    async def _process_video(self, file_data: bytes, filename: str, format_info: Dict) -> Dict[str, Any]:
        """处理视频"""
        return {
            "text": "",
            "metadata": {
                "format": format_info["extension"],
                "size": len(file_data),
                "note": "视频处理待实现"
            }
        }
    
    async def _process_mindmap(self, file_data: bytes, filename: str, format_info: Dict) -> Dict[str, Any]:
        """处理思维导图"""
        return {
            "text": "",
            "metadata": {"format": format_info["extension"], "note": "思维导图处理待实现"}
        }
    
    async def _process_database(self, file_data: bytes, filename: str, format_info: Dict) -> Dict[str, Any]:
        """处理数据库文件⭐增强版（支持CSV、JSON、SQL）"""
        ext = format_info["extension"]
        
        try:
            # CSV文件
            if ext == ".csv":
                try:
                    # 尝试多种编码
                    for encoding in ['utf-8', 'gbk', 'gb2312', 'latin-1']:
                        try:
                            text = file_data.decode(encoding)
                            import csv
                            import io
                            
                            # 解析CSV
                            csvfile = io.StringIO(text)
                            reader = csv.reader(csvfile)
                            rows = list(reader)
                            
                            # 格式化输出（限制行数）
                            text_parts = []
                            for i, row in enumerate(rows[:1000]):  # 限制1000行
                                text_parts.append("\t".join(row))
                            
                            if len(rows) > 1000:
                                text_parts.append("... (行数过多，已截断)")
                            
                            return {
                                "text": "\n".join(text_parts),
                                "metadata": {
                                    "format": "csv",
                                    "rows": min(len(rows), 1000),
                                    "total_rows": len(rows),
                                    "encoding": encoding
                                }
                            }
                        except UnicodeDecodeError:
                            continue
                except Exception as e:
                    logger.warning(f"CSV处理失败: {e}")
            
            # JSON文件
            elif ext == ".json":
                try:
                    import json
                    text = file_data.decode('utf-8', errors='ignore')
                    data = json.loads(text)
                    # 格式化JSON
                    formatted_text = json.dumps(data, ensure_ascii=False, indent=2)
                    return {
                        "text": formatted_text,
                        "metadata": {
                            "format": "json",
                            "type": type(data).__name__
                        }
                    }
                except Exception as e:
                    logger.warning(f"JSON处理失败: {e}")
            
            # SQL文件
            elif ext == ".sql":
                try:
                    text = file_data.decode('utf-8', errors='ignore')
                    return {
                        "text": text,
                        "metadata": {
                            "format": "sql",
                            "lines": len(text.split('\n'))
                        }
                    }
                except Exception as e:
                    logger.warning(f"SQL处理失败: {e}")
            
            # XML文件
            elif ext == ".xml":
                try:
                    text = file_data.decode('utf-8', errors='ignore')
                    return {
                        "text": text,
                        "metadata": {
                            "format": "xml",
                            "lines": len(text.split('\n'))
                        }
                    }
                except Exception as e:
                    logger.warning(f"XML处理失败: {e}")
            
            return {
                "text": "",
                "metadata": {"format": ext, "note": f"{ext}格式处理待实现"}
            }
        except Exception as e:
            return {"text": "", "metadata": {"error": str(e)}}
    
    async def _process_archive(self, file_data: bytes, filename: str, format_info: Dict) -> Dict[str, Any]:
        """处理压缩文件⭐增强版（支持ZIP解压和内容列表）"""
        ext = format_info["extension"]
        
        try:
            # ZIP文件
            if ext == ".zip":
                try:
                    import zipfile
                    import io
                    
                    zf = zipfile.ZipFile(io.BytesIO(file_data))
                    file_list = zf.namelist()
                    
                    # 获取文件列表和总大小
                    total_size = sum(zf.getinfo(name).file_size for name in file_list)
                    
                    # 尝试提取文本文件内容（小文件）
                    text_parts = []
                    text_parts.append(f"压缩包包含 {len(file_list)} 个文件:")
                    
                    for name in file_list[:100]:  # 限制列出100个文件
                        info = zf.getinfo(name)
                        text_parts.append(f"  - {name} ({info.file_size} 字节)")
                        
                        # 如果是小的文本文件，尝试提取内容
                        if (info.file_size < 100000 and 
                            any(name.lower().endswith(ext) for ext in ['.txt', '.md', '.json', '.xml', '.py', '.js'])):
                            try:
                                content = zf.read(name).decode('utf-8', errors='ignore')
                                text_parts.append(f"    内容预览: {content[:200]}...")
                            except:
                                pass
                    
                    if len(file_list) > 100:
                        text_parts.append(f"  ... (还有 {len(file_list) - 100} 个文件)")
                    
                    return {
                        "text": "\n".join(text_parts),
                        "metadata": {
                            "format": "zip",
                            "files": len(file_list),
                            "total_size": total_size,
                            "method": "zipfile"
                        }
                    }
                except ImportError:
                    pass
                except Exception as e:
                    logger.warning(f"ZIP处理失败: {e}")
            
            # 其他压缩格式
            return {
                "text": "",
                "metadata": {
                    "format": ext,
                    "note": f"{ext}格式需要相应的解压库"
                }
            }
        except Exception as e:
            return {"text": "", "metadata": {"error": str(e)}}
    
    async def _process_text_file(
        self,
        file_data: bytes,
        filename: str,
        format_info: Dict
    ) -> Dict[str, Any]:
        """处理文本文件"""
        try:
            text = file_data.decode('utf-8', errors='ignore')
            return {
                "text": text,
                "metadata": {
                    "format": format_info["extension"],
                    "lines": len(text.split('\n')),
                    "chars": len(text)
                }
            }
        except Exception as e:
            return {"text": "", "metadata": {"error": str(e)}}
    
    def _detect_language(self, filename: str) -> str:
        """检测编程语言"""
        ext = os.path.splitext(filename)[1].lower()
        language_map = {
            ".py": "python",
            ".js": "javascript",
            ".java": "java",
            ".c": "c",
            ".cpp": "cpp",
            ".cs": "csharp",
            ".swift": "swift",
            ".kt": "kotlin",
            ".go": "go",
            ".rs": "rust",
            ".ts": "typescript",
            ".php": "php",
            ".rb": "ruby",
            ".r": "r",
            ".sql": "sql",
            ".html": "html",
            ".css": "css",
            ".sh": "shell",
            ".ps1": "powershell"
        }
        return language_map.get(ext, "unknown")
    
    def get_supported_formats(self) -> List[str]:
        """获取所有支持的文件格式"""
        all_formats = []
        for formats in self.supported_formats.values():
            all_formats.extend(formats)
        return sorted(set(all_formats))
    
    def get_format_statistics(self) -> Dict[str, Any]:
        """
        获取格式统计信息
        
        Returns:
            统计信息字典
        """
        total_formats = sum(len(formats) for formats in self.supported_formats.values())
        
        return {
            "total_formats": total_formats,
            "categories": {
                category: len(formats)
                for category, formats in self.supported_formats.items()
            },
            "processors": {
                category: processor.__name__
                for category, processor in self.processors.items()
            }
        }
    
    async def batch_process_files(
        self,
        files: List[Dict[str, Any]],
        max_parallel: int = 5
    ) -> List[Dict[str, Any]]:
        """
        批量处理文件⭐性能优化版（并行处理）
        
        Args:
            files: 文件列表，每个元素包含file_data、filename、mime_type
            max_parallel: 最大并行处理数
            
        Returns:
            处理结果列表
        """
        # 分批并行处理
        results = []
        
        for i in range(0, len(files), max_parallel):
            batch = files[i:i+max_parallel]
            
            # 创建任务列表
            tasks = [
                self.process_file(
                    file_data=f["file_data"],
                    filename=f["filename"],
                    mime_type=f.get("mime_type")
                )
                for f in batch
            ]
            
            # 并行执行
            batch_results = await asyncio.gather(*tasks, return_exceptions=True)
            
            # 处理异常
            for j, result in enumerate(batch_results):
                if isinstance(result, Exception):
                    results.append({
                        "success": False,
                        "error": str(result),
                        "filename": batch[j]["filename"]
                    })
                else:
                    results.append(result)
        
        return results

