"""
文件生成服务
支持生成Word/Excel/PPT/PDF等文件
"""

from typing import Dict, List, Optional, Any
from datetime import datetime
import io
import json

class FileGenerationService:
    """
    文件生成服务
    
    功能：
    1. Word文档生成
    2. Excel表格生成
    3. PPT演示文稿生成
    4. PDF文档生成
    5. 图片生成
    6. 模板管理
    7. 与RAG知识库集成（自动保存生成的文件内容）
    8. 与内容创作模块集成（导出生成的内容）
    """
    
    def __init__(self, rag_service=None, content_service=None):
        """
        初始化文件生成服务
        
        Args:
            rag_service: RAG服务适配器（可选）
            content_service: 内容创作服务（可选）
        """
        self.templates = {}
        self.use_libraries = {
            "word": False,  # python-docx
            "excel": True,  # openpyxl (已在ERP中使用)
            "ppt": False,   # python-pptx
            "pdf": False    # reportlab/weasyprint
        }
        self.rag_service = rag_service  # RAG服务（用于自动保存）
        self.content_service = content_service  # 内容服务（用于导出）
        
    async def generate_word(
        self,
        content: str,
        template: Optional[str] = None,
        output_path: Optional[str] = None,
        title: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        生成Word文档
        
        Args:
            content: 内容（Markdown或HTML格式）
            template: 模板名称
            output_path: 输出路径
            title: 文档标题
            
        Returns:
            包含文件数据和元数据的字典
        """
        try:
            # 检查是否安装了python-docx
            try:
                from docx import Document
                from docx.shared import Inches, Pt
                from docx.enum.text import WD_ALIGN_PARAGRAPH
                self.use_libraries["word"] = True
            except ImportError:
                self.use_libraries["word"] = False
            
            if self.use_libraries["word"]:
                # 使用python-docx生成Word文档
                doc = Document()
                
                # 添加标题
                if title:
                    heading = doc.add_heading(title, 0)
                    heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
                
                # 解析Markdown内容（简单实现）
                lines = content.split('\n')
                for line in lines:
                    line = line.strip()
                    if not line:
                        doc.add_paragraph()
                        continue
                    
                    # 处理标题
                    if line.startswith('# '):
                        doc.add_heading(line[2:], level=1)
                    elif line.startswith('## '):
                        doc.add_heading(line[3:], level=2)
                    elif line.startswith('### '):
                        doc.add_heading(line[4:], level=3)
                    # 处理列表
                    elif line.startswith('- ') or line.startswith('* '):
                        doc.add_paragraph(line[2:], style='List Bullet')
                    elif line.startswith('1. ') or line.startswith('2. '):
                        doc.add_paragraph(line[3:], style='List Number')
                    else:
                        doc.add_paragraph(line)
                
                # 保存到内存
                buffer = io.BytesIO()
                doc.save(buffer)
                buffer.seek(0)
                file_data = buffer.read()
                
                result = {
                    "success": True,
                    "file_data": file_data,
                    "format": "docx",
                    "filename": output_path or f"document_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx",
                    "size": len(file_data),
                    "timestamp": datetime.now().isoformat()
                }
                
                # 自动保存到RAG知识库
                if self.rag_service:
                    try:
                        await self._save_to_rag(content, title or "Word文档", "docx", result)
                    except Exception as e:
                        result["rag_save_warning"] = f"保存到RAG失败: {str(e)}"
                
                return result
            else:
                # 备用方案：返回HTML格式（浏览器可以打开）
                html_content = self._markdown_to_html(content, title)
                return {
                    "success": True,
                    "file_data": html_content.encode('utf-8'),
                    "format": "html",
                    "filename": output_path or f"document_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html",
                    "note": "python-docx未安装，返回HTML格式",
                    "timestamp": datetime.now().isoformat()
                }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "timestamp": datetime.now().isoformat()
            }
    
    def _markdown_to_html(self, markdown: str, title: Optional[str] = None) -> str:
        """简单的Markdown转HTML"""
        html = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>{title or '文档'}</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 40px; line-height: 1.6; }}
        h1 {{ color: #333; border-bottom: 2px solid #333; }}
        h2 {{ color: #555; margin-top: 30px; }}
        h3 {{ color: #777; }}
        ul, ol {{ margin-left: 20px; }}
        code {{ background: #f4f4f4; padding: 2px 5px; border-radius: 3px; }}
    </style>
</head>
<body>
"""
        if title:
            html += f"<h1>{title}</h1>\n"
        
        lines = markdown.split('\n')
        for line in lines:
            if line.startswith('# '):
                html += f"<h1>{line[2:]}</h1>\n"
            elif line.startswith('## '):
                html += f"<h2>{line[3:]}</h2>\n"
            elif line.startswith('### '):
                html += f"<h3>{line[4:]}</h3>\n"
            elif line.startswith('- ') or line.startswith('* '):
                html += f"<li>{line[2:]}</li>\n"
            else:
                html += f"<p>{line}</p>\n"
        
        html += """
</body>
</html>
"""
        return html
    
    async def generate_excel(
        self,
        data: List[List[Any]],
        headers: Optional[List[str]] = None,
        output_path: Optional[str] = None,
        sheet_name: str = "Sheet1"
    ) -> Dict[str, Any]:
        """
        生成Excel表格
        
        Args:
            data: 数据（二维列表）
            headers: 表头
            output_path: 输出路径
            sheet_name: 工作表名称
            
        Returns:
            包含文件数据和元数据的字典
        """
        try:
            # 使用openpyxl（已在ERP中使用）
            try:
                from openpyxl import Workbook
                from openpyxl.styles import Font, Alignment
                self.use_libraries["excel"] = True
            except ImportError:
                self.use_libraries["excel"] = False
            
            if self.use_libraries["excel"]:
                wb = Workbook()
                ws = wb.active
                ws.title = sheet_name
                
                # 添加表头
                if headers:
                    for col_idx, header in enumerate(headers, 1):
                        cell = ws.cell(row=1, column=col_idx, value=header)
                        cell.font = Font(bold=True)
                        cell.alignment = Alignment(horizontal='center')
                
                # 添加数据
                start_row = 2 if headers else 1
                for row_idx, row_data in enumerate(data, start_row):
                    for col_idx, value in enumerate(row_data, 1):
                        ws.cell(row=row_idx, column=col_idx, value=value)
                
                # 自动调整列宽
                for col in ws.columns:
                    max_length = 0
                    col_letter = col[0].column_letter
                    for cell in col:
                        try:
                            if len(str(cell.value)) > max_length:
                                max_length = len(str(cell.value))
                        except:
                            pass
                    adjusted_width = min(max_length + 2, 50)
                    ws.column_dimensions[col_letter].width = adjusted_width
                
                # 保存到内存
                buffer = io.BytesIO()
                wb.save(buffer)
                buffer.seek(0)
                file_data = buffer.read()
                
                result = {
                    "success": True,
                    "file_data": file_data,
                    "format": "xlsx",
                    "filename": output_path or f"excel_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx",
                    "size": len(file_data),
                    "rows": len(data) + (1 if headers else 0),
                    "columns": len(headers) if headers else (len(data[0]) if data else 0),
                    "timestamp": datetime.now().isoformat()
                }
                
                # 自动保存到RAG知识库（将表格数据转换为文本）
                if self.rag_service:
                    try:
                        table_text = self._table_to_text(headers, data)
                        await self._save_to_rag(table_text, f"Excel表格-{sheet_name}", "xlsx", result)
                    except Exception as e:
                        result["rag_save_warning"] = f"保存到RAG失败: {str(e)}"
                
                return result
            else:
                # 备用方案：返回CSV格式
                csv_lines = []
                if headers:
                    csv_lines.append(",".join(str(h) for h in headers))
                for row in data:
                    csv_lines.append(",".join(str(v) for v in row))
                csv_content = "\n".join(csv_lines)
                
                return {
                    "success": True,
                    "file_data": csv_content.encode('utf-8-sig'),  # UTF-8 BOM for Excel
                    "format": "csv",
                    "filename": output_path or f"excel_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv",
                    "note": "openpyxl未安装，返回CSV格式",
                    "timestamp": datetime.now().isoformat()
                }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "timestamp": datetime.now().isoformat()
            }
    
    async def generate_ppt(
        self,
        slides: List[Dict[str, Any]],
        template: Optional[str] = None,
        output_path: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        生成PPT演示文稿
        
        Args:
            slides: 幻灯片列表，每个slide包含：
                - title: 标题
                - content: 内容（文本或列表）
                - bullet_points: 要点列表（可选）
            template: 模板名称
            output_path: 输出路径
            
        Returns:
            包含文件数据和元数据的字典
        """
        try:
            # 检查是否安装了python-pptx
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt
                from pptx.enum.text import PP_ALIGN
                self.use_libraries["ppt"] = True
            except ImportError:
                self.use_libraries["ppt"] = False
            
            if self.use_libraries["ppt"]:
                # 使用python-pptx生成PPT
                prs = Presentation()
                
                # 设置幻灯片尺寸（16:9）
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(5.625)
                
                # 为每个slide创建幻灯片
                for slide_data in slides:
                    # 使用标题和内容布局
                    slide_layout = prs.slide_layouts[1]  # 标题和内容布局
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # 设置标题
                    title = slide.shapes.title
                    title.text = slide_data.get("title", "无标题")
                    
                    # 设置内容
                    content = slide_data.get("content", "")
                    bullet_points = slide_data.get("bullet_points", [])
                    
                    # 获取内容占位符
                    if len(slide.placeholders) > 1:
                        content_placeholder = slide.placeholders[1]
                        tf = content_placeholder.text_frame
                        tf.text = content
                        
                        # 添加要点
                        if bullet_points:
                            for point in bullet_points:
                                p = tf.add_paragraph()
                                p.text = point
                                p.level = 0
                
                # 保存到内存
                buffer = io.BytesIO()
                prs.save(buffer)
                buffer.seek(0)
                file_data = buffer.read()
                
                result = {
                    "success": True,
                    "file_data": file_data,
                    "format": "pptx",
                    "filename": output_path or f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                    "size": len(file_data),
                    "slides_count": len(slides),
                    "timestamp": datetime.now().isoformat()
                }
                
                # 自动保存到RAG知识库（将幻灯片内容转换为文本）
                if self.rag_service:
                    try:
                        slides_text = self._slides_to_text(slides)
                        await self._save_to_rag(slides_text, "PPT演示文稿", "pptx", result)
                    except Exception as e:
                        result["rag_save_warning"] = f"保存到RAG失败: {str(e)}"
                
                return result
            else:
                # 备用方案：返回HTML格式（浏览器可以打开）
                html_content = self._slides_to_html(slides)
                return {
                    "success": True,
                    "file_data": html_content.encode('utf-8'),
                    "format": "html",
                    "filename": output_path or f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html",
                    "note": "python-pptx未安装，返回HTML格式",
                    "slides_count": len(slides),
                    "timestamp": datetime.now().isoformat()
                }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "timestamp": datetime.now().isoformat()
            }
    
    def _slides_to_html(self, slides: List[Dict[str, Any]]) -> str:
        """将幻灯片转换为HTML格式"""
        html = """<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>演示文稿</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 40px; }
        .slide { page-break-after: always; margin-bottom: 50px; }
        h1 { color: #333; border-bottom: 2px solid #333; }
        h2 { color: #555; margin-top: 20px; }
        ul { margin-left: 20px; }
        .slide-number { position: absolute; bottom: 10px; right: 10px; color: #999; }
    </style>
</head>
<body>
"""
        for idx, slide_data in enumerate(slides, 1):
            html += f'<div class="slide">\n'
            html += f'<h1>{slide_data.get("title", "无标题")}</h1>\n'
            
            content = slide_data.get("content", "")
            if content:
                html += f'<p>{content}</p>\n'
            
            bullet_points = slide_data.get("bullet_points", [])
            if bullet_points:
                html += '<ul>\n'
                for point in bullet_points:
                    html += f'<li>{point}</li>\n'
                html += '</ul>\n'
            
            html += f'<div class="slide-number">第 {idx} 页</div>\n'
            html += '</div>\n'
        
        html += """
</body>
</html>
"""
        return html
    
    async def generate_pdf(
        self,
        content: str,
        template: Optional[str] = None,
        output_path: Optional[str] = None,
        title: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        生成PDF文档
        
        Args:
            content: 内容（HTML或Markdown格式）
            template: 模板名称
            output_path: 输出路径
            title: 文档标题
            
        Returns:
            包含文件数据和元数据的字典
        """
        try:
            # 尝试使用weasyprint（推荐，支持HTML转PDF）
            try:
                from weasyprint import HTML
                self.use_libraries["pdf"] = True
                pdf_method = "weasyprint"
            except ImportError:
                # 尝试使用reportlab
                try:
                    from reportlab.lib.pagesizes import letter, A4
                    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
                    from reportlab.lib.styles import getSampleStyleSheet
                    self.use_libraries["pdf"] = True
                    pdf_method = "reportlab"
                except ImportError:
                    self.use_libraries["pdf"] = False
                    pdf_method = None
            
            if self.use_libraries["pdf"] and pdf_method == "weasyprint":
                # 使用weasyprint生成PDF
                html_content = self._markdown_to_html(content, title)
                pdf_bytes = HTML(string=html_content).write_pdf()
                
                result = {
                    "success": True,
                    "file_data": pdf_bytes,
                    "format": "pdf",
                    "filename": output_path or f"document_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf",
                    "size": len(pdf_bytes),
                    "method": "weasyprint",
                    "timestamp": datetime.now().isoformat()
                }
                
                # 自动保存到RAG知识库
                if self.rag_service:
                    try:
                        await self._save_to_rag(content, title or "PDF文档", "pdf", result)
                    except Exception as e:
                        result["rag_save_warning"] = f"保存到RAG失败: {str(e)}"
                
                return result
            elif self.use_libraries["pdf"] and pdf_method == "reportlab":
                # 使用reportlab生成PDF
                buffer = io.BytesIO()
                doc = SimpleDocTemplate(buffer, pagesize=A4)
                styles = getSampleStyleSheet()
                story = []
                
                if title:
                    story.append(Paragraph(title, styles['Title']))
                    story.append(Spacer(1, 12))
                
                # 解析内容
                lines = content.split('\n')
                for line in lines:
                    line = line.strip()
                    if not line:
                        story.append(Spacer(1, 6))
                        continue
                    story.append(Paragraph(line, styles['Normal']))
                
                doc.build(story)
                buffer.seek(0)
                pdf_bytes = buffer.read()
                
                result = {
                    "success": True,
                    "file_data": pdf_bytes,
                    "format": "pdf",
                    "filename": output_path or f"document_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf",
                    "size": len(pdf_bytes),
                    "method": "reportlab",
                    "timestamp": datetime.now().isoformat()
                }
                
                # 自动保存到RAG知识库
                if self.rag_service:
                    try:
                        await self._save_to_rag(content, title or "PDF文档", "pdf", result)
                    except Exception as e:
                        result["rag_save_warning"] = f"保存到RAG失败: {str(e)}"
                
                return result
            else:
                # 备用方案：返回HTML格式（浏览器可以打印为PDF）
                html_content = self._markdown_to_html(content, title)
                return {
                    "success": True,
                    "file_data": html_content.encode('utf-8'),
                    "format": "html",
                    "filename": output_path or f"document_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html",
                    "note": "PDF库未安装，返回HTML格式（可在浏览器中打印为PDF）",
                    "timestamp": datetime.now().isoformat()
                }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "timestamp": datetime.now().isoformat()
            }
    
    async def generate_image(
        self,
        content: str,
        width: int = 800,
        height: int = 600,
        output_path: Optional[str] = None
    ) -> bytes:
        """
        生成图片
        
        Args:
            content: 内容（文本或HTML）
            width: 宽度
            height: 高度
            output_path: 输出路径
            
        Returns:
            图片字节流
        """
        # TODO: 使用PIL/Pillow生成图片
        
        return b""
    
    def _table_to_text(self, headers: Optional[List[str]], data: List[List[Any]]) -> str:
        """将表格数据转换为文本格式"""
        text_lines = []
        if headers:
            text_lines.append(" | ".join(str(h) for h in headers))
            text_lines.append("-" * (sum(len(str(h)) for h in headers) + len(headers) * 3))
        for row in data:
            text_lines.append(" | ".join(str(v) for v in row))
        return "\n".join(text_lines)
    
    def _slides_to_text(self, slides: List[Dict[str, Any]]) -> str:
        """将幻灯片内容转换为文本格式"""
        text_lines = []
        for idx, slide_data in enumerate(slides, 1):
            text_lines.append(f"\n=== 幻灯片 {idx}: {slide_data.get('title', '无标题')} ===\n")
            content = slide_data.get("content", "")
            if content:
                text_lines.append(content)
            bullet_points = slide_data.get("bullet_points", [])
            if bullet_points:
                for point in bullet_points:
                    text_lines.append(f"  • {point}")
        return "\n".join(text_lines)
    
    async def _save_to_rag(
        self,
        content: str,
        title: str,
        file_format: str,
        file_metadata: Dict[str, Any]
    ):
        """
        将生成的文件内容保存到RAG知识库
        
        Args:
            content: 文件文本内容
            title: 文档标题
            file_format: 文件格式（docx/xlsx/pptx/pdf）
            file_metadata: 文件元数据
        """
        if not self.rag_service:
            return
        
        try:
            # 构造文档数据
            doc_text = f"标题: {title}\n格式: {file_format}\n生成时间: {file_metadata.get('timestamp', '')}\n\n内容:\n{content}"
            
            # 调用RAG服务保存
            if hasattr(self.rag_service, 'ingest_text'):
                await self.rag_service.ingest_text(
                    text=doc_text,
                    doc_id=f"file_gen_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
                    metadata={
                        "source": "file_generation",
                        "file_format": file_format,
                        "filename": file_metadata.get("filename", ""),
                        "size": file_metadata.get("size", 0),
                        "generated_at": file_metadata.get("timestamp", "")
                    },
                    save_index=True
                )
            elif hasattr(self.rag_service, 'add_document'):
                # 兼容其他RAG服务接口
                await self.rag_service.add_document(
                    content=doc_text,
                    title=title,
                    metadata={
                        "source": "file_generation",
                        "file_format": file_format,
                        **file_metadata
                    }
                )
        except Exception as e:
            # 不抛出异常，只记录警告
            import logging
            logger = logging.getLogger(__name__)
            logger.warning(f"保存文件到RAG失败: {str(e)}")
    
    async def export_content_to_file(
        self,
        content_data: Dict[str, Any],
        file_type: str = "docx",
        output_path: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        从内容创作模块导出内容为文件
        
        Args:
            content_data: 内容数据（包含title, content, tags等）
            file_type: 文件类型（docx/xlsx/pptx/pdf）
            output_path: 输出路径
            
        Returns:
            文件生成结果
        """
        title = content_data.get("title", "内容")
        content = content_data.get("content", "")
        
        if file_type == "docx":
            return await self.generate_word(content, title=title, output_path=output_path)
        elif file_type == "pdf":
            return await self.generate_pdf(content, title=title, output_path=output_path)
        elif file_type == "pptx":
            # 将内容转换为幻灯片格式
            slides = [{
                "title": title,
                "content": content,
                "bullet_points": content_data.get("tags", [])
            }]
            return await self.generate_ppt(slides, output_path=output_path)
        else:
            return {
                "success": False,
                "error": f"不支持的文件类型: {file_type}",
                "timestamp": datetime.now().isoformat()
            }

